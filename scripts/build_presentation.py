"""
研究内容をまとめたPowerPointを生成する。
既存の results/*.png をそのまま流用し、追加の図表生成は行わない。
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

RESULTS = os.path.expanduser("~/projects/fugaku-power-forecast/results")
OUT_PATH = os.path.expanduser("~/projects/fugaku-power-forecast/fugaku_power_forecast_presentation.pptx")

INK = RGBColor(0x17, 0x18, 0x1A)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0xB3, 0x73, 0x1F)
ACCENT_SOFT = RGBColor(0xF2, 0xE2, 0xC8)
PAPER = RGBColor(0xF7, 0xF6, 0xF2)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE1, 0xE0, 0xD9)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
RED = RGBColor(0xE3, 0x49, 0x48)
GOOD = RGBColor(0x0C, 0xA3, 0x0C)

FONT_JP = "Meiryo"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font=FONT_JP, line_spacing=1.15, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15.5, color=INK2, font=FONT_JP,
                 line_spacing=1.25, space_after=10, marker="—  "):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = marker + item
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = font
    return box


def add_kicker(slide, text, top=Inches(0.42)):
    add_text(slide, Inches(0.6), top, Inches(6), Inches(0.4), text, size=13, color=ACCENT,
              bold=True, font=FONT_JP)


def add_pagenum(slide, n):
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35), str(n),
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer_rule(slide, y=Inches(0.95)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(12.13), Pt(1.4))
    line.fill.solid(); line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def section(slide_idx, kicker, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, PAPER)
    add_kicker(slide, kicker)
    add_text(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.9), title, size=30, bold=True, color=INK)
    y = Inches(1.55)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6), subtitle, size=16, color=INK2)
        y = Inches(2.1)
    add_footer_rule(slide, y=Inches(1.42) if not subtitle else Inches(2.0))
    add_pagenum(slide, slide_idx)
    return slide


def image_slide(slide_idx, kicker, title, img_path, bullets=None, img_top=None, img_width=None, caption=None):
    slide = section(slide_idx, kicker, title)
    content_top = Inches(1.75)
    if bullets:
        add_bullets(slide, Inches(0.6), content_top, Inches(12.1), Inches(1.3), bullets, size=15)
        content_top = Inches(3.05)
    if img_top:
        content_top = img_top
    from PIL import Image
    im = Image.open(img_path)
    iw, ih = im.size
    max_w = img_width or Inches(11.5)
    max_h = Inches(7.5) - content_top - Inches(0.35)
    scale = min(max_w / iw, max_h / ih)
    w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
    left = Emu(int((SLIDE_W - w) / 2))
    slide.shapes.add_picture(img_path, left, content_top, width=w, height=h)
    if caption:
        add_text(slide, Inches(0.6), Inches(7.05), Inches(11), Inches(0.35), caption, size=10.5, color=MUTED)
    return slide


def add_table(slide, left, top, width, col_widths, headers, rows, header_fill=SURFACE.__class__(0x17,0x18,0x1A),
              font_size=12.5, row_h=Inches(0.42)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_h * n_rows
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table
    for c, cw in enumerate(col_widths):
        table.columns[c].width = cw
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = FONT_JP
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(8); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = SURFACE if r % 2 == 0 else RGBColor(0xF1, 0xEF, 0xE8)
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = INK2 if c > 0 else INK
                run.font.name = FONT_JP
                if c == 0:
                    run.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(8); cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
    return table


def add_stat(slide, left, top, width, value, label, value_color=ACCENT):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.25))
    box.fill.solid(); box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Pt(14); tf.margin_top = Pt(10); tf.margin_right = Pt(10)
    p0 = tf.paragraphs[0]
    p0.text = value
    p0.alignment = PP_ALIGN.LEFT
    for run in p0.runs:
        run.font.size = Pt(26); run.font.bold = True; run.font.color.rgb = value_color; run.font.name = "Consolas"
    p1 = tf.add_paragraph()
    p1.text = label
    p1.alignment = PP_ALIGN.LEFT
    for run in p1.runs:
        run.font.size = Pt(11.5); run.font.color.rgb = MUTED; run.font.name = FONT_JP


# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(BLANK)
add_bg(slide, INK)
add_text(slide, Inches(0.9), Inches(2.5), Inches(4), Inches(0.4), "研究発表資料 ・ 全5ステップ完了",
         size=13, color=RGBColor(0xF0, 0xB8, 0x62), bold=True)
add_text(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(2.0),
         "富岳の消費電力を予測する", size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(slide, Inches(0.9), Inches(4.05), Inches(11), Inches(1.0),
         "事前学習済み時系列モデル vs 稼働シグナルを用いた自作MLモデル\n"
         "―― レジーム遷移期における予測の頑健性と、その実務的価値",
         size=17, color=RGBColor(0xC3, 0xC2, 0xB7), line_spacing=1.4)
add_text(slide, Inches(0.9), Inches(6.7), Inches(8), Inches(0.4),
         "Fugaku Power Forecast Project", size=12, color=MUTED)

# ============================================================
# Slide 2: Agenda
# ============================================================
slide = section(2, "AGENDA", "本日お話しする5つのステップ")
items = [
    ("01", "データパイプライン", "2000万件のジョブログを1時間単位の消費電力系列に再構成"),
    ("02", "7モデル比較", "事前学習済みモデル(Chronos) vs 3アーキテクチャ×{電力のみ/covariate込み}"),
    ("03", "訂正", "「遷移期でこそcovariateが効く」という当初の結論を、サンプルを増やして検証し直す"),
    ("04", "早期警戒シグナル", "SHAP・Attentionで「どの指標を監視すべきか」を特定"),
    ("05", "意思決定シミュレーション", "予測を冷却容量の引当判断に接続し、運用コストへの効果を定量化"),
]
y = 1.85
for num, title, desc in items:
    add_text(slide, Inches(0.6), Inches(y), Inches(0.9), Inches(0.5), num, size=22, bold=True, color=ACCENT, font="Consolas")
    add_text(slide, Inches(1.5), Inches(y - 0.03), Inches(3.3), Inches(0.5), title, size=17, bold=True, color=INK)
    add_text(slide, Inches(4.9), Inches(y), Inches(7.6), Inches(0.5), desc, size=13, color=INK2)
    y += 1.0

# ============================================================
# Slide 3: 背景・問い
# ============================================================
slide = section(3, "背景", "なぜこの研究をするのか")
add_bullets(slide, Inches(0.6), Inches(1.85), Inches(11.8), Inches(3.5), [
    "スーパーコンピュータの消費電力を事前に予測できれば、冷却設備や電力バジェットの計画を最適化できる",
    "時系列予測には大きく2つのアプローチがある：",
    "　①既に大量データで学習済みの「時系列基盤モデル」をそのまま使う（zero-shot）",
    "　②手元のデータで「稼働状況を表す追加情報(covariate)」まで使って自作モデルを一から学習する",
    "本研究の核心的な問い：稼働状況の情報は、特にワークロードが急変する「レジーム遷移期」において",
    "　②を①より優位にできるのか？　そしてその予測は実務上どれだけの価値を持つのか？",
], size=16.5, line_spacing=1.5, space_after=14)

# ============================================================
# Slide 4: Step1 データパイプライン
# ============================================================
image_slide(4, "STEP 01", "データパイプライン：ジョブログ→時間単位の電力系列",
            os.path.join(RESULTS, "hourly_timeseries_overview.png"),
            bullets=["約2000万件のジョブ記録を、実行時間の重なりで按分し、28,084時間分の系列に再構成",
                     "妥当性検証：再構成した稼働ノード数のピーク158,886 ≈ 富岳の実ノード数158,976（誤差0.06%）"],
            img_width=Inches(10.3))

# ============================================================
# Slide 5: Step2 モデル設計
# ============================================================
slide = section(5, "STEP 02", "7モデルを同一条件で比較する設計")
add_table(slide, Inches(0.6), Inches(1.85), Inches(12.1),
          [Inches(3.6), Inches(4.25), Inches(4.25)],
          ["アーキテクチャ", "②電力のみ", "③電力+covariate"],
          [["LightGBM（勾配ブースティング木）", "②-LightGBM", "③-LightGBM"],
           ["LSTM（RNN）", "②-LSTM", "③-LSTM"],
           ["TFT-lite（Attention系）", "②-TFT", "③-TFT"],
           ["Chronos-bolt（事前学習済み、zero-shot）", "①として参照点に使用（学習なし）", "―"]],
          row_h=Inches(0.55))
add_bullets(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(2.2), [
    "①vs②：入力情報量を揃えた上での「事前学習知識 vs ゼロからの学習」の効果",
    "②vs③：同一アーキテクチャでのcovariate情報の効果",
    "test期間中は全モデルとも重み凍結。Chronosとの比較を公平にするため、参照できる過去情報は共通で直近720時間に統一",
], size=14.5, line_spacing=1.4)

# ============================================================
# Slide 6: Step2 結果 MASE比較
# ============================================================
image_slide(6, "STEP 02 ・ 結果", "事前学習済みモデルが、ゼロから学習した全モデルを上回った",
            os.path.join(RESULTS, "step2_overall_mase.png"),
            bullets=["MASE（低いほど良い）：Chronos 0.654 ＜ 各アーキテクチャの②③（0.74〜0.79）"],
            img_width=Inches(9.0))

# ============================================================
# Slide 7: Step2 covariate効果の分岐
# ============================================================
slide = section(7, "STEP 02 ・ 結果", "covariateの効果はアーキテクチャ次第だった")
add_table(slide, Inches(0.6), Inches(1.85), Inches(12.1),
          [Inches(3.0), Inches(3.03), Inches(3.03), Inches(3.03)],
          ["architecture", "短期・安定期", "短期・遷移期", "中期・遷移期"],
          [["LightGBM", "+1.6%", "+5.3%", "+2.8%"],
           ["LSTM", "-1.8%", "+4.5%", "+3.1%"],
           ["TFT-lite", "+0.3%", "-5.3%", "-2.7%"]],
          row_h=Inches(0.55))
add_bullets(slide, Inches(0.6), Inches(4.25), Inches(12), Inches(2.6), [
    "LightGBM・LSTM：covariateの効果は安定期より遷移期で明確に大きい → 当初の仮説を支持",
    "TFT-liteだけは逆に遷移期でcovariateが悪化要因に。パラメータ数が増える一方で学習データ量(1,796件)は"
    "変わらず、過学習した可能性",
    "→ 複数アーキテクチャで比較したからこそ見えた「covariateの恩恵はアーキテクチャ依存」という発見",
], size=15, line_spacing=1.5, space_after=12)

# ============================================================
# Slide 8: Step3 訂正
# ============================================================
slide = section(8, "STEP 03", "訂正：サンプルを増やしたら、目玉の結果が消えた")
add_bullets(slide, Inches(0.6), Inches(1.85), Inches(11.9), Inches(1.0), [
    "Step2の「遷移期」は単一の変化点(2022-11-08)のみに基づき、サンプル数が少なすぎた(n=673/174)",
    "test期間内部に絞って高感度な変化点検出を行うと5点が見つかり、サンプル数はn=3,365/870に増加",
], size=15, line_spacing=1.4, space_after=10)
add_table(slide, Inches(0.6), Inches(3.05), Inches(12.1),
          [Inches(3.0), Inches(4.55), Inches(4.55)],
          ["architecture", "短期・遷移期 (Step2, n=673)", "短期・遷移期 (Step3精緻化, n=3365)"],
          [["LightGBM", "+5.3%", "-0.3%"],
           ["LSTM", "+4.5%", "-0.9%"],
           ["TFT-lite", "-5.3%", "-0.9%"]],
          row_h=Inches(0.55))
add_text(slide, Inches(0.6), Inches(5.45), Inches(12), Inches(1.4),
         "一方、安定期での③の優位はStep2・Step3で一貫（LightGBM +1.6%→+2.3%）。\n"
         "→「covariateは遷移期でこそ効く」という当初の主張は再現せず、研究として重要な自己訂正となった。",
         size=15.5, color=INK, line_spacing=1.5)

# ============================================================
# Slide 9: Step3 検知遅延
# ============================================================
image_slide(9, "STEP 03 ・ 別の角度から", "「平均誤差」では見えなかった効果が、「回復の速さ」では見えた",
            os.path.join(RESULTS, "step3_detection_lag_curves.png"),
            bullets=["遷移直後、1時間先予測の誤差が基準値の1.3倍以内に戻るまでの日数：",
                     "LSTMのみcovariate追加で 7.0日 → 4.6日 に短縮（-34%）。LightGBM/TFT-liteは差なし"],
            img_width=Inches(11.3))

# ============================================================
# Slide 10: Step4 早期警戒シグナル
# ============================================================
slide = section(10, "STEP 04", "見るべきシグナルは「到着数」より「待ち時間」")
add_bullets(slide, Inches(0.6), Inches(1.85), Inches(11.9), Inches(1.9), [
    "LightGBMのSHAP分析とTFT-liteのAttention(Variable Selection)重み ―― 2つの独立した解釈手法が一致",
    "covariate全体の寄与度は安定期/遷移期でほぼ変化しない（LightGBM: 16.5%→16.5%）",
    "しかし「どのcovariateが効くか」は遷移期で明確にシフトする：",
], size=15, line_spacing=1.4, space_after=8)
add_table(slide, Inches(0.6), Inches(3.75), Inches(12.1),
          [Inches(4.5), Inches(3.8), Inches(3.8)],
          ["covariate", "遷移期/安定期 重要度比", "解釈"],
          [["待ち時間の急増z-score(168h)", "1.70倍", "最も遷移期で重要度が上がる"],
           ["compute-bound比率(168h平均)", "1.44倍", "ワークロード種別の週内シフト"],
           ["新規ユーザ出現z-score(720h)", "1.24倍", "新規プロジェクト流入の兆候"],
           ["ジョブ到着数(168h平均・単独最大)", "0.91倍", "強い予測子だが遷移の予兆としては相対的に鈍い"]],
          row_h=Inches(0.55), font_size=13)

# ============================================================
# Slide 11: Step5 枠組み
# ============================================================
slide = section(11, "STEP 05", "予測を「容量引当」の意思決定に接続する")
add_bullets(slide, Inches(0.6), Inches(1.85), Inches(11.9), Inches(2.6), [
    "冷却容量／電力バジェットをどれだけ事前に確保するかという判断を想定",
    "少なく確保→急上昇時に緊急対応（高コスト）／多く確保→余剰容量が無駄（低コスト）",
    "このコスト最小化問題の最適解は、数学的に「分位点予測」そのもの",
    "→ Step2で計算済みの分位点予測(q=0.75, q=0.90)をそのまま流用でき、追加学習は不要だった",
], size=16, line_spacing=1.55, space_after=14)
add_text(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.6),
         "cost(y, C) = r × max(y − C, 0) + max(C − y, 0)　　(r = 過小引当コスト ÷ 過大引当コスト)",
         size=15, color=ACCENT, font="Consolas", bold=True)
add_text(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(1.3),
         "比較対象：static（train期間の固定分位点で常に確保）／ seasonal-naive（前日の値＋誤差マージン）／ 7モデル",
         size=14, color=INK2)

# ============================================================
# Slide 12: Step5 結果
# ============================================================
slide = section(12, "STEP 05 ・ 結果", "予測を使うと、引当コストが最大34.5%下がる")
add_stat(slide, Inches(0.6), Inches(1.85), Inches(3.85), "-34.5%", "最良モデルのコスト対 seasonal-naive（高リスク回避シナリオ）")
add_stat(slide, Inches(4.6), Inches(1.85), Inches(3.85), "-17.6%", "最良モデルのコスト対 static（予測なし）安全マージン")
add_stat(slide, Inches(8.6), Inches(1.85), Inches(3.85), "0.75 / 0.90", "使用した分位点。Step2で計算済み、追加学習なし")
add_bullets(slide, Inches(0.6), Inches(3.4), Inches(11.9), Inches(2.8), [
    "7モデル全てが両ベースラインを上回った",
    "意外な発見：seasonal-naiveはstaticより悪化することがある（5.39M vs 4.28M）。富岳の電力は"
    "メンテナンスで急激にゼロ近くまで落ち込むため、素朴な予測はこの急変に弱い",
    "→「質の低い予測は何もしないより悪い」。7モデルをきちんと比較検証したこと自体に実務的意義がある",
    "レジーム別に見ても、LightGBMのcovariate効果は安定期でのみコスト削減（-2.7%）とStep3の結論と整合",
], size=15, line_spacing=1.45, space_after=10)

# ============================================================
# Slide 13: 総括
# ============================================================
slide = section(13, "総括", "5つの結論")
concl = [
    ("① 事前学習 vs ②ゼロ学習", "Chronos zero-shotが一貫して最良。富岳15ヶ月分のデータでは、大規模時系列コーパスでの事前学習に対抗できなかった。"),
    ("② vs ③ covariateの効果", "アーキテクチャ依存。LightGBMは一貫して(小幅に)効き、LSTMは回復速度でのみ効き、TFT-liteは過学習で逆効果。"),
    ("研究上の自己訂正", "「遷移期でこそcovariateが効く」という初期の主張は、サンプルを増やすと再現しなかった。"),
    ("早期警戒シグナル", "到着率そのものより、待ち時間(キュー長)の急増z-scoreの方が遷移の予兆として感度が高い。"),
    ("実務的価値", "動的な容量引当は静的方式に対し最大34.5%のコスト削減。ただし質の低い予測は逆効果になりうる。"),
]
y = 1.85
for i, (title, desc) in enumerate(concl):
    add_text(slide, Inches(0.6), Inches(y), Inches(0.5), Inches(0.4), str(i + 1), size=18, bold=True, color=ACCENT, font="Consolas")
    add_text(slide, Inches(1.15), Inches(y), Inches(3.1), Inches(0.85), title, size=13.5, bold=True, color=INK, line_spacing=1.2)
    add_text(slide, Inches(4.4), Inches(y), Inches(8.1), Inches(0.85), desc, size=13, color=INK2, line_spacing=1.3)
    y += 1.02

# ============================================================
# Slide 14: 今後の課題
# ============================================================
slide = section(14, "今後の課題", "残された発展的テーマ")
add_bullets(slide, Inches(0.6), Inches(1.9), Inches(11.9), Inches(4), [
    "分位点予測のキャリブレーション補正（conformal prediction） ― 7モデルとも分位点がやや過小に出ている",
    "Chronosのfine-tuning ― zero-shotで既に最強だが、富岳データで追加学習すればさらに伸びるか",
    "変化点検出のハイパラ感度分析 ― pen値・遷移期window幅(±14日)の選び方への結果の依存度を検証",
    "LSTMの解釈性分析 ― LightGBM(SHAP)・TFT(Attention)と揃えるため、Integrated Gradients等の適用",
    "TFT-liteの自己注意重みの時間方向分析 ― 「どの過去時刻が効いたか」の可視化",
], size=16, line_spacing=1.6, space_after=16)
add_text(slide, Inches(0.6), Inches(6.5), Inches(11.5), Inches(0.6),
         "詳細：~/projects/fugaku-power-forecast/README.md", size=12.5, color=MUTED, font="Consolas")

prs.save(OUT_PATH)
print(f"saved: {OUT_PATH}")
print(f"slides: {len(prs.slides.slides) if hasattr(prs.slides,'slides') else len(prs.slides._sldIdLst)}")
